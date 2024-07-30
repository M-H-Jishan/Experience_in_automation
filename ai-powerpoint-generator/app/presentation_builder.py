from pptx import Presentation
from pptx.util import Inches, Pt
from io import BytesIO
from app.slide_generator import generate_slide_titles, generate_slide_content
from app.image_generator import generate_image

def create_presentation(topic, num_slides, template):
    prs = Presentation(f"templates/{template}")
    titles = generate_slide_titles(topic, num_slides)
    
    for title in titles:
        slide_layout = prs.slide_layouts[1]  # Using a layout with a title and content
        slide = prs.slides.add_slide(slide_layout)
        
        # Set the slide title
        slide.shapes.title.text = title
        
        # Generate and set the slide content
        content = generate_slide_content(topic, title)
        content_placeholder = slide.placeholders[1]
        text_frame = content_placeholder.text_frame
        text_frame.clear()
        for point in content:
            p = text_frame.add_paragraph()
            p.text = point
            p.level = 0
        
        # Generate and add an image
        image_prompt = f"Create an image related to '{title}' for a presentation on '{topic}'"
        image_data = generate_image(image_prompt)
        img_stream = BytesIO(image_data)
        slide.shapes.add_picture(img_stream, Inches(5), Inches(2), width=Inches(4))
    
    return prs